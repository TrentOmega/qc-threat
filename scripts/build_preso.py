#!/usr/bin/env python3
"""
Generate the Quantum Threat to Bitcoin presentation.

Produces a .pptx then optionally converts to .odp via LibreOffice.
Theme: black (#000) / orange (#F7931A) / amber-gold (#F4C27A).

Usage:
    python scripts/build_preso.py            # build .pptx + .odp
    python scripts/build_preso.py --no-odp   # .pptx only
"""

import subprocess
import sys
import re
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── paths ────────────────────────────────────────────────────────────────────

ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "assets" / "images"
CHART = ROOT / "assets" / "charts"
LOGO = IMG / "bitcoinbushbash-logo.png"
OUT = ROOT / "quantum-threat-bitcoin-bushbash.pptx"
PRESO_MD = ROOT / "qc-threat-pres.md"

# ── palette ──────────────────────────────────────────────────────────────────

BLACK   = RGBColor(0x00, 0x00, 0x00)
NEAR    = RGBColor(0x0D, 0x0D, 0x0D)  # card bg
EDGE    = RGBColor(0x2A, 0x2A, 0x2A)  # card border
GRID    = RGBColor(0x1A, 0x1A, 0x1A)
OFF_W   = RGBColor(0xF5, 0xF0, 0xE8)  # warm off-white body text
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE  = RGBColor(0xF7, 0x93, 0x1A)
GOLD    = RGBColor(0xF4, 0xC2, 0x7A)
RED     = RGBColor(0xE4, 0x57, 0x2E)
MUTED   = RGBColor(0x88, 0x88, 0x88)

H_BG      = "#000000"
H_TEXT    = "#F5F0E8"
H_ACCENT  = "#F7931A"
H_GOLD    = "#F4C27A"
H_RED     = "#E4572E"
H_MUTED   = "#888888"

FONT = "Liberation Sans"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── GRI appendix A.4 ────────────────────────────────────────────────────────

GRI_HORIZONS = ["5 yr", "10 yr", "15 yr", "20 yr", "30 yr"]
GRI_ROWS = [
    ("Extremely unlikely (<1%)",            [18, 4, 1, 0, 0], 1,   0),
    ("Very unlikely (<5%)",                 [5, 11, 3, 1, 1], 5,   1),
    ("Unlikely (<30%)",                     [6,  7, 7, 2, 0], 30,  5),
    ("Neither likely nor unlikely (~50%)",  [2,  5,10,10, 3], 70, 30),
    ("Likely (>70%)",                       [1,  3, 6, 8,12], 95, 70),
    ("Very likely (>95%)",                  [0,  2, 4, 7,10], 99, 95),
    ("Extremely likely (>99%)",             [0,  0, 1, 4, 6],100, 99),
]
N_EXPERTS = 32


def gri_derived():
    years = [5, 10, 15, 20, 30]
    opt, pes = [], []
    for i in range(5):
        o = sum(r[1][i] * r[2] for r in GRI_ROWS) / N_EXPERTS
        p = sum(r[1][i] * r[3] for r in GRI_ROWS) / N_EXPERTS
        opt.append(o)
        pes.append(p)
    return years, opt, pes


# ═══════════════════════════════════════════════════════════════════════════════
#  CHART GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def _ax_dark(ax):
    ax.set_facecolor(H_BG)
    ax.tick_params(colors=H_TEXT, labelsize=10)
    for s in ax.spines.values():
        s.set_color("#333")


def chart_exposure():
    plt.style.use("dark_background")
    names = ["P2PK", "P2TR\nkey-path", "Reused\nP2PKH", "Fresh\nP2WPKH", "P2SH/P2WSH\n(unspent)"]
    vals  = [95, 82, 68, 34, 18]
    cols  = [H_ACCENT, H_GOLD, H_RED, "#C9772B", "#666"]

    fig, ax = plt.subplots(figsize=(10, 5), dpi=150)
    bars = ax.bar(names, vals, color=cols, width=0.62)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width()/2, v + 2.5, str(v),
                ha="center", color=H_TEXT, fontsize=12, weight="bold")
    ax.set_ylabel("Relative risk", color=H_TEXT, fontsize=12)
    ax.set_ylim(0, 108)
    ax.set_title("Long-Exposure Quantum Risk by Script Type", color=H_TEXT, fontsize=16, pad=10)
    ax.grid(axis="y", alpha=0.12)
    _ax_dark(ax)
    fig.patch.set_facecolor(H_BG)
    fig.tight_layout()
    path = CHART / "exposure-risk.png"
    fig.savefig(path, facecolor=H_BG)
    plt.close(fig)
    return path


def chart_gri_curve():
    plt.style.use("dark_background")
    yrs, opt, pes = gri_derived()

    fig, ax = plt.subplots(figsize=(10, 5), dpi=150)
    ax.plot(yrs, opt, color=H_ACCENT, lw=3, marker="o", ms=7, label="Optimistic")
    ax.plot(yrs, pes, color=H_GOLD,   lw=3, marker="o", ms=7, label="Pessimistic")
    ax.fill_between(yrs, pes, opt, color="#333", alpha=0.35)
    for x, y in zip(yrs, opt):
        ax.annotate(f"{y:.0f}%", (x, y), textcoords="offset points",
                    xytext=(0, 10), ha="center", color=H_ACCENT, fontsize=10, weight="bold")
    for x, y in zip(yrs, pes):
        ax.annotate(f"{y:.0f}%", (x, y), textcoords="offset points",
                    xytext=(0, -15), ha="center", color=H_GOLD, fontsize=10)
    ax.set_title("GRI 2024 — Derived CRQC Probability", color=H_TEXT, fontsize=16, pad=10)
    ax.set_xlabel("Horizon (years)", color=H_TEXT, fontsize=12)
    ax.set_ylabel("Average probability %", color=H_TEXT, fontsize=12)
    ax.set_ylim(0, 100)
    ax.set_xticks(yrs)
    ax.grid(alpha=0.12)
    _ax_dark(ax)
    leg = ax.legend(frameon=False, fontsize=11)
    for t in leg.get_texts():
        t.set_color(H_TEXT)
    fig.patch.set_facecolor(H_BG)
    fig.tight_layout()
    path = CHART / "gri-curve.png"
    fig.savefig(path, facecolor=H_BG)
    plt.close(fig)
    return path


def chart_gri_table():
    plt.style.use("dark_background")
    fig, ax = plt.subplots(figsize=(12, 5.5), dpi=160)
    fig.patch.set_facecolor(H_BG)
    ax.set_facecolor(H_BG)
    ax.axis("off")

    cell_text = [r[1] for r in GRI_ROWS]
    row_labels = [r[0] for r in GRI_ROWS]

    tbl = ax.table(cellText=cell_text, rowLabels=row_labels, colLabels=GRI_HORIZONS,
                   loc="center", cellLoc="center", rowLoc="center")
    tbl.scale(1, 1.55)
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor("#444")
        cell.set_linewidth(0.7)
        cell.get_text().set_color(H_TEXT)
        cell.get_text().set_fontsize(10)
        if r == 0:
            cell.set_facecolor("#1a1a1a")
            cell.get_text().set_color(H_ACCENT)
            cell.get_text().set_weight("bold")
        elif c == -1:
            cell.set_facecolor("#111")
            cell.get_text().set_color(H_GOLD)
            cell.get_text().set_fontsize(9)
        else:
            cell.set_facecolor("#080808")
    ax.set_title("GRI 2024 — Expert Response Counts (N = 32)", color=H_TEXT, fontsize=16, pad=14)
    fig.tight_layout()
    path = CHART / "gri-table.png"
    fig.savefig(path, facecolor=H_BG, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)
    return path


def chart_migration():
    plt.style.use("dark_background")
    yrs = list(range(2026, 2035))
    std    = [2,  8, 16, 28, 41, 55, 67, 77, 85]
    urgent = [3, 13, 27, 43, 58, 71, 82, 90, 95]
    risk   = [1,  3,  7, 13, 23, 35, 49, 62, 72]

    fig, ax = plt.subplots(figsize=(10, 5), dpi=150)
    ax.plot(yrs, urgent, color=H_ACCENT, lw=3, marker="s", ms=5, label="Urgent rollout")
    ax.plot(yrs, std,    color=H_GOLD,   lw=3, marker="o", ms=5, label="Standard rollout")
    ax.plot(yrs, risk,   color=H_RED,    lw=2.5, ls="--", label="CRQC risk signal")
    ax.fill_between(yrs, risk, urgent, alpha=0.06, color=H_ACCENT)
    ax.set_title("Migration Coverage vs Q-Day Risk", color=H_TEXT, fontsize=16, pad=10)
    ax.set_ylabel("% coverage / risk index", color=H_TEXT, fontsize=12)
    ax.set_ylim(0, 100)
    ax.set_xticks(yrs)
    ax.grid(alpha=0.12)
    _ax_dark(ax)
    leg = ax.legend(frameon=False, fontsize=11)
    for t in leg.get_texts():
        t.set_color(H_TEXT)
    fig.patch.set_facecolor(H_BG)
    fig.tight_layout()
    path = CHART / "migration.png"
    fig.savefig(path, facecolor=H_BG)
    plt.close(fig)
    return path


def chart_witness():
    """Horizontal bar: P2MR witness sizes vs P2TR."""
    plt.style.use("dark_background")
    labels = ["P2TR key-path", "P2MR depth 0", "P2MR depth 1",
              "P2MR depth 3", "P2TR script\ndepth 1"]
    sizes  = [66, 103, 135, 199, 167]
    cols   = [H_RED, H_ACCENT, H_ACCENT, H_ACCENT, H_GOLD]

    fig, ax = plt.subplots(figsize=(9, 4.5), dpi=150)
    bars = ax.barh(labels, sizes, color=cols, height=0.55)
    for b, v in zip(bars, sizes):
        ax.text(v + 3, b.get_y() + b.get_height()/2, f"{v} B",
                va="center", color=H_TEXT, fontsize=11, weight="bold")
    ax.set_xlim(0, 230)
    ax.set_title("Witness Size Comparison (bytes)", color=H_TEXT, fontsize=15, pad=10)
    ax.grid(axis="x", alpha=0.12)
    _ax_dark(ax)
    fig.patch.set_facecolor(H_BG)
    fig.tight_layout()
    path = CHART / "witness-cmp.png"
    fig.savefig(path, facecolor=H_BG)
    plt.close(fig)
    return path


def chart_merkle_tree():
    """Diagram of P2MR Merkle tree structure."""
    plt.style.use("dark_background")
    fig, ax = plt.subplots(figsize=(10, 5.5), dpi=150)
    fig.patch.set_facecolor(H_BG)
    ax.set_facecolor(H_BG)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 6)
    ax.axis("off")

    def node(x, y, txt, fc, ec, tc, fs=13):
        ax.text(x, y, txt, ha="center", va="center", fontsize=fs, color=tc, weight="bold",
                bbox=dict(boxstyle="round,pad=0.35", facecolor=fc, edgecolor=ec, lw=1.5))

    def edge(x1, y1, x2, y2):
        ax.plot([x1, x2], [y1, y2], color="#555", lw=1.3)

    node(5, 5.1, "Merkle Root  (32 B)", "#1a1a1a", H_ACCENT, H_ACCENT, 15)
    node(2.8, 3.5, "TapBranch", "#111", H_GOLD, H_GOLD)
    node(7.2, 3.5, "TapBranch", "#111", H_GOLD, H_GOLD)
    for x, lbl in [(1.1, "Leaf A\ntapscript"), (4.5, "Leaf B\ntapscript"),
                   (5.9, "Leaf C\nfuture PQC"), (8.9, "Leaf D\ntapscript")]:
        node(x, 1.7, lbl, "#0a0a0a", "#555", H_TEXT, 11)

    edge(5, 4.7, 2.8, 3.9)
    edge(5, 4.7, 7.2, 3.9)
    edge(2.8, 3.1, 1.1, 2.15)
    edge(2.8, 3.1, 4.5, 2.15)
    edge(7.2, 3.1, 5.9, 2.15)
    edge(7.2, 3.1, 8.9, 2.15)

    ax.text(5, 0.55, "scriptPubKey:  OP_2  OP_PUSHBYTES_32  <root>",
            ha="center", fontsize=11, color=H_MUTED, style="italic")
    ax.text(5, 0.1, "No internal key.  No key-path spend.  Address prefix: bc1z",
            ha="center", fontsize=10, color=H_RED)

    path = CHART / "merkle-tree.png"
    fig.savefig(path, facecolor=H_BG, bbox_inches="tight", pad_inches=0.15)
    plt.close(fig)
    return path


def generate_charts():
    CHART.mkdir(parents=True, exist_ok=True)
    return {
        "exposure":  chart_exposure(),
        "gri_curve": chart_gri_curve(),
        "gri_table": chart_gri_table(),
        "migration": chart_migration(),
        "witness":   chart_witness(),
        "merkle":    chart_merkle_tree(),
    }


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE PRIMITIVES
# ═══════════════════════════════════════════════════════════════════════════════

def _fill_bg(slide, prs):
    r = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                               Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = BLACK
    r.line.fill.background()


def _logo(slide):
    slide.shapes.add_picture(str(LOGO), Inches(11.7), Inches(0.12), width=Inches(1.48))


def _footer(slide):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(4), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.text = "bitcoinbushbash.info"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.font.name = FONT


def _accent_bar(slide, x=0.7, y=1.35, w=4.5):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(0.055))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()


def _title(slide, text, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.25), Inches(10.5), Inches(1.1))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    _accent_bar(slide)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(10.5), Inches(0.45))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(17)
        sp.font.color.rgb = GOLD
        sp.font.name = FONT


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def slide_base(prs, title_text, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(s, prs)
    _logo(s)
    _footer(s)
    _title(s, title_text, subtitle)
    return s


def _bullets(slide, items, x=0.7, y=2.1, w=6.0, h=4.5, size=26, color=None, bold=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color or OFF_W
        p.font.name = FONT
        p.space_after = Pt(12)


def _two_col(slide, l_head, l_items, r_head, r_items, y=2.05):
    for cx, head in [(0.7, l_head), (7.0, r_head)]:
        tb = slide.shapes.add_textbox(Inches(cx), Inches(y), Inches(5.5), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = head
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        p.font.name = FONT

    for cx, items, w in [(0.7, l_items, 5.8), (7.0, r_items, 5.3)]:
        tb = slide.shapes.add_textbox(Inches(cx), Inches(y + 0.55), Inches(w), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, txt in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = txt
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = OFF_W
            p.font.name = FONT
            p.space_after = Pt(9)


def _img(slide, path, x=6.2, y=2.0, w=6.2, h=4.8):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def _img_full(slide, path, x=0.5, y=1.9, w=12.3, h=5.1):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def _card(slide, x, y, w, h, big, label):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = NEAR
    box.line.color.rgb = EDGE
    box.line.width = Pt(1)

    nt = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.25), Inches(w - 0.2), Inches(0.8))
    p = nt.text_frame.paragraphs[0]
    p.text = big
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    dt = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.15), Inches(w - 0.2), Inches(0.8))
    p2 = dt.text_frame.paragraphs[0]
    p2.text = label
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = OFF_W
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER


def _cards_row(slide, items, y=2.3):
    n = len(items)
    gap = 0.2
    total = 12.0
    w = (total - gap * (n - 1)) / n
    for i, (big, label) in enumerate(items):
        _card(slide, 0.65 + i * (w + gap), y, w, 2.2, big, label)


def _callout(slide, text, x=0.7, y=5.5, w=11.8, color=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(17)
    p.font.color.rgb = color or GOLD
    p.font.name = FONT


def load_reference_sections(md_path: Path):
    text = md_path.read_text(encoding="utf-8")
    lines = text.splitlines()

    in_refs = False
    section = None
    sections = []
    current = None
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped == "## References":
            in_refs = True
            i += 1
            continue

        if in_refs and stripped.startswith("## ") and stripped != "## References":
            break

        if not in_refs:
            i += 1
            continue

        if stripped.startswith("### "):
            section = {"title": stripped[4:].strip(), "entries": []}
            sections.append(section)
            current = None
            i += 1
            continue

        match = re.match(r"^(\d+)\.\s+\[(.+?)\]\((https?://[^)]+)\)", stripped)
        if match and section is not None:
            current = {
                "number": match.group(1),
                "title": match.group(2),
                "url": match.group(3),
                "desc": "",
            }
            section["entries"].append(current)
            i += 1
            continue

        if current is not None and stripped:
            current["desc"] = f"{current['desc']} {stripped}".strip()

        i += 1

    return [s for s in sections if s["entries"]]


def add_reference_slides(prs):
    sections = load_reference_sections(PRESO_MD)
    refs_by_number = {}
    for section in sections:
        for entry in section["entries"]:
            refs_by_number[int(entry["number"])] = entry

    page_plan = [
        ("What is QC", [1, 2, 3]),
        ("What is the threat", [4, 5, 6, 7]),
        ("Bitcoin revision", [8, 9, 10, 11, 12]),
        ("QC timeline", [13, 14, 15]),
        ("Technical problem and fixes", [16, 17]),
        (None, [18, 19, 20, 21, 22]),
        (None, [23, 24, 25]),
        (None, [26, 27]),
        ("Social problem and fixes", [28, 29, 30, 31]),
        ("Other uses of QC", [32, 33]),
        (None, [34, 35, 36]),
    ]

    for idx, (section_title, ref_numbers) in enumerate(page_plan, start=1):
        title = "References" if len(page_plan) == 1 else f"References ({idx}/{len(page_plan)})"
        s = slide_base(prs, title)

        body = s.shapes.add_textbox(Inches(0.7), Inches(1.95), Inches(12.0), Inches(5.0))
        tf = body.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

        first = True
        if section_title:
            p = tf.paragraphs[0]
            p.text = section_title
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ORANGE
            p.font.name = FONT
            p.space_before = Pt(8)
            p.space_after = Pt(4)
            first = False

        for ref_num in ref_numbers:
            entry = refs_by_number[ref_num]
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            title_run = p.add_run()
            title_run.text = f"{entry['number']}. {entry['title']}"
            title_run.font.size = Pt(12)
            title_run.font.bold = True
            title_run.font.color.rgb = OFF_W
            title_run.font.name = FONT
            title_run.hyperlink.address = entry["url"]
            p.space_before = Pt(0)
            p.space_after = Pt(1)

            if entry["desc"]:
                p2 = tf.add_paragraph()
                desc_run = p2.add_run()
                desc_run.text = entry["desc"]
                desc_run.font.size = Pt(12)
                desc_run.font.bold = False
                desc_run.font.color.rgb = MUTED
                desc_run.font.name = FONT
                p2.space_before = Pt(0)
                p2.space_after = Pt(6)

        _notes(s, "")


# ═══════════════════════════════════════════════════════════════════════════════
#  DECK
# ═══════════════════════════════════════════════════════════════════════════════

def build_deck(charts: dict) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── 1  Title ─────────────────────────────────────────────────────────────
    s = slide_base(prs, "The Quantum Threat to Bitcoin",
                   "What's at risk, what's being built, and what you can do")
    _bullets(s, [
        "Overview the threat",
        "What we can do",
        "What you can do",
    ], y=2.5, size=30)
    _notes(s, "30 min including questions. Bitcoin beginner-advanced, QC none-intermediate.\n"
              "Frame as a migration race, not overnight collapse.")

    # ── 2  Talk map ──────────────────────────────────────────────────────────
    s = slide_base(prs, "Talk map")
    _bullets(s, [
        "1. What is quantum computing",
        "2. The QC threat to Bitcoin",
        "3. Revision: sigs, hashes, mempool",
        "4. When is it coming (GRI 2024)",
        "5. Two big problems",
        "6. BIP-360 deep dive  (the big mitigation)",
        "7. Legacy coins dilemma",
        "8. What you can do",
        "9. Mythbusters & perspective",
    ], y=2.05, size=24)
    _notes(s, "Anchor slide. BIP-360 gets 5 dedicated slides.")

    # ── 3  What is QC ────────────────────────────────────────────────────────
    s = slide_base(prs, "What is quantum computing?",
                   "Only what matters for Bitcoin's security model")
    _two_col(s,
        "Core ideas",
        ["Superposition: 0 and 1 at once",
         "Entanglement: correlated states",
         "Interference: amplify right answers"],
        "Road to a CRQC",
        ["Physical qubits are noisy",
         "Error correction burns many physical qubits per logical qubit",
         "Logical qubits + runtime = can it run Shor fast enough?"])
    _notes(s, "CRQC = Cryptographically Relevant Quantum Computer.\n"
              "Current largest factored number by QC is 15. Far away, but trajectory matters.")

    # ── 4  Threat classes ────────────────────────────────────────────────────
    s = slide_base(prs, "The QC threat to Bitcoin",
                   "Two algorithms, two threat profiles")
    _two_col(s,
        "Shor's algorithm  (signatures)",
        ["Breaks ECDSA and Schnorr",
         "Derives private key from public key",
         "Exponential speedup",
         "Primary Bitcoin threat"],
        "Grover's algorithm  (hashing)",
        ["Targets SHA-256 / RIPEMD-160",
         "Quadratic speedup only",
         "256-bit -> ~128-bit effective",
         "Changes mining economics, not collapse"])
    _notes(s, "Shor: exponential advantage, breaks DLP. The real worry.\n"
              "Grover: quadratic only. Not a collapse event, a cost-curve shift.")

    # ── 5  Revision ──────────────────────────────────────────────────────────
    s = slide_base(prs, "Revision",
                   "Digital signatures, hashing, and the mempool")
    _bullets(s, [
        "Signatures: prove you own the coins, authorize spend",
        "Hashing: mining proof-of-work + compress/conceal public keys",
        "Public key hidden by hash until first spend (P2PKH/P2SH/P2WSH)",
        "Public key exposed on receive (P2PK, P2TR key-path)",
        "Mempool window: tx waits ~10 min before confirmation",
        "When your public key is visible = your quantum risk window",
    ], y=2.1, size=24)
    _notes(s, "Key insight: spending behaviour determines exposure.\n"
              "Address reuse = key exposed on-chain forever.\n"
              "P2TR key-path = key exposed at receive time (like 2010-era P2PK).")

    # ── 6  Long vs short exposure ────────────────────────────────────────────
    s = slide_base(prs, "Attack types: long vs short exposure")
    _img_full(s, charts["exposure"])
    _notes(s, "Long exposure: key visible on-chain indefinitely. Attacker has unlimited time.\n"
              "  P2PK/P2TR key-path: exposed at receive.\n"
              "  Reused P2PKH: exposed after first spend.\n"
              "Short exposure: mempool race after key reveal. Needs faster QC. Secondary threat.")

    # ── 7  GRI headline ──────────────────────────────────────────────────────
    s = slide_base(prs, "When is it coming?",
                   "GRI Quantum Threat Timeline Report 2024  —  32 experts surveyed")
    _cards_row(s, [
        ("~14%", "Optimistic avg\nwithin 5 years"),
        ("~34%", "Optimistic avg\nwithin 10 years"),
        ("~19%", "Pessimistic avg\nwithin 10 years"),
    ])
    _callout(s, "Planning signal, not prediction. Uncertainty concentrates in the 10-15 year window.", y=5.2)
    _notes(s, "Values from GRI 2024 key points. Weighted averages of 32 expert responses.")

    # ── 8  GRI table ─────────────────────────────────────────────────────────
    s = slide_base(prs, "GRI expert response distribution",
                   "Raw counts — Appendix A.4")
    _img_full(s, charts["gri_table"])
    _notes(s, "Each column sums to 32. Shows how spread the expert opinion really is.\n"
              "18 of 32 say extremely unlikely within 5 years. But 3 say likely or above.")

    # ── 9  GRI derived curve ─────────────────────────────────────────────────
    s = slide_base(prs, "GRI derived probability curve",
                   "Optimistic vs pessimistic probability assignments")
    _img_full(s, charts["gri_curve"])
    _notes(s, "Derived from Appendix A.4 raw counts * probability bin assignments.\n"
              "Key: ~34% optimistic at 10 yr, ~62% at 15 yr. The window is tightening.")

    # ── 10  Government timelines ─────────────────────────────────────────────
    s = slide_base(prs, "Government policy timelines",
                   "Regulatory pressure already moving")
    _cards_row(s, [
        ("2030", "CNSA 2.0\nPQC mandate begins"),
        ("2033", "CNSA 2.0\nfull browser/OS PQC"),
        ("2035", "NIST IR 8547\ndisallow ECC"),
    ])
    _callout(s, "Bitcoin doesn't answer to regulators — but enterprise/institutional users will demand PQC.", y=5.2)
    _notes(s, "CNSA 2.0: Commercial National Security Algorithm Suite.\n"
              "NIST IR 8547 disallows ECC in US federal gov after 2035.\n"
              "These drive enterprise Bitcoin infrastructure timelines.")

    # ── 11  Two big problems ─────────────────────────────────────────────────
    s = slide_base(prs, "The two big Bitcoin problems")
    _bullets(s, [
        "1.  Technical:  add post-quantum cryptography to Bitcoin",
        "2.  Social:  what happens to legacy coins that don't upgrade",
    ], y=2.5, w=11.0, size=32)
    _callout(s, "Both must be solved together. Technical alone is not enough.")
    _notes(s, "Pivot into solution space. These are coupled problems.")

    # ── 12  PQC families ─────────────────────────────────────────────────────
    s = slide_base(prs, "Post-quantum signature options",
                   "5 families, different tradeoffs")
    _two_col(s,
        "NIST standardised",
        ["ML-DSA (Dilithium) — lattice-based",
         "SLH-DSA (SPHINCS+) — hash-based",
         "Conservative, well-studied",
         "Compact sigs, larger public keys"],
        "Other families",
        ["Code-based (McEliece) — huge keys",
         "Multivariate (Rainbow — broken)",
         "Isogeny (SIKE — broken)",
         "Active research, real risks"])
    _callout(s, "BIP-360 does NOT commit to any PQC algorithm — it builds the infrastructure for them.")
    _notes(s, "ML-DSA and SLH-DSA are the NIST PQC standards (Aug 2024).\n"
              "Some candidates already broken: SIKE, Rainbow. Caution warranted.")

    # ── 13  PQC impacts ──────────────────────────────────────────────────────
    s = slide_base(prs, "PQC operational impacts",
                   "Bigger signatures ripple through the entire stack")
    _bullets(s, [
        "Transaction size and fee increases",
        "Block space pressure and mempool dynamics",
        "Verification and signing time changes",
        "Lightning channel updates: every state = PQC signature",
        "Multisig: N signatures * PQC size",
        "Exchanges and custody: key management overhaul",
    ], y=2.1, size=24)
    _notes(s, "PQC sigs are much larger than ECDSA/Schnorr (64B -> potentially kilobytes).\n"
              "Affects on-chain, Lightning, custody, hardware wallets.")

    # ═══ BIP-360 DEEP DIVE (slides 14-18) ════════════════════════════════════

    # ── 14  BIP-360 overview ─────────────────────────────────────────────────
    s = slide_base(prs, "BIP-360: Pay-to-Merkle-Root",
                   "The key quantum mitigation path for Bitcoin")
    _cards_row(s, [
        ("P2MR", "New output type\nSegWit version 2"),
        ("bc1z...", "New address\nprefix (Bech32m)"),
        ("Soft fork", "Backward\ncompatible"),
    ])
    _bullets(s, [
        "Removes key-path spend — no public key exposed on-chain",
        "Reuses battle-tested tapscript (BIPs 340/341/342)",
        "Creates safe upgrade path for future PQC opcodes",
    ], y=5.0, w=11.5, size=19)
    _notes(s, "BIP-360 by Hunter Beast. Status: Draft.\n"
              "Core idea: commit to Merkle root of script tree, not a tweaked public key.\n"
              "Eliminates long-exposure quantum attack surface entirely.\n"
              "BitMEX endorses this approach: removes Taproot key-path to match P2PKH quantum safety.")

    # ── 15  BIP-360 how it works ─────────────────────────────────────────────
    s = slide_base(prs, "BIP-360: Merkle tree structure",
                   "How P2MR commits to scripts without exposing keys")
    _img_full(s, charts["merkle"])
    _notes(s, "Output: OP_2 OP_PUSHBYTES_32 <merkle_root>\n"
              "Unlike P2TR: no internal key, no key tweak. Commits only to TapBranch hash.\n"
              "Witness: stack elements + leaf script + control block (NO public key in control block).\n"
              "Control block: 1 + 32*m bytes (m = tree depth, max 128).\n"
              "Parity bit always 1 (no key to compute parity from).\n"
              "Unused script paths remain private. Tapscript compatible.")

    # ── 16  BIP-360 witness sizes ────────────────────────────────────────────
    s = slide_base(prs, "BIP-360: witness size tradeoffs",
                   "Larger than P2TR key-path — smaller than P2TR script-path")
    _img(s, charts["witness"], x=0.5, y=1.95, w=7.2, h=4.8)
    _bullets(s, [
        "Depth 0 (single leaf): 103 B",
        "vs P2TR key-path: 66 B",
        "+37 B base overhead",
        "+32 B per tree level",
        "",
        "Always 32 B smaller than",
        "equivalent P2TR script-path",
        "(no internal key in control block)",
    ], x=7.9, y=2.1, w=4.8, size=19)
    _notes(s, "P2MR witness formula: 37 + 32*m bytes more than P2TR key-path.\n"
              "But always 32 bytes LESS than equivalent P2TR script-path spend.\n"
              "Because P2MR omits the internal public key from control block.\n"
              "Depth 0: 103B. Depth 1: 135B. Depth 3: 199B.")

    # ── 17  BIP-360 design rationale ─────────────────────────────────────────
    s = slide_base(prs, "BIP-360: design principles",
                   "Minimum change, maximum future optionality")
    _two_col(s,
        "Design goals (from BIP text)",
        ["Minimize changes to the network",
         "Reuse existing Bitcoin code and workflows",
         "Preserve compatibility where possible",
         "Safest path for PQC signature integration",
         "Enable gradual, iterative upgrades"],
        "Key technical decisions",
        ["No key-path spend (eliminates ECC surface)",
         "No internal key tweak (Merkle root only)",
         "OP_SUCCESSx path for future PQC opcodes",
         "No specific PQC algorithm mandated",
         "Multiple PQC sigs possible for redundancy"])
    _notes(s, "Verbatim from BIP-360: 'Minimize changes to the network.'\n"
              "'Create the safest possible path for post-quantum signature integrations.'\n"
              "'Facilitate gradual integration of quantum resistant features carried out iteratively.'\n"
              "Future PQC via OP_SUCCESSx opcode updates for PQC OP_CHECKSIG.")

    # ── 18  BIP-360 protection scope ─────────────────────────────────────────
    s = slide_base(prs, "BIP-360: what it protects (and doesn't)",
                   "Long-exposure resistance now — short-exposure resistance later")
    _two_col(s,
        "Protected: long exposure",
        ["Public keys never appear on-chain",
         "Funds at rest are quantum-safe",
         "Unused script paths stay private",
         "Existing tapscript programs work as-is"],
        "Not yet protected: short exposure",
        ["Mempool window still uses ECDSA/Schnorr",
         "Key revealed during broadcast-to-confirm",
         "Needs future PQC signature opcodes",
         "P2MR creates the path for these"])
    _callout(s, "P2MR is stage 1.  PQC signatures via OP_SUCCESSx are stage 2.")
    _notes(s, "From BIP-360: 'P2MR outputs are only resistant to long exposure attacks.'\n"
              "'Protection against short exposure attacks may require future post-quantum signatures.'\n"
              "BitMEX proposes dual tapleaf: quantum-safe AND quantum-vulnerable options per address.")

    # ── 19  Migration timeline ───────────────────────────────────────────────
    s = slide_base(prs, "Implementation timeline",
                   "We win if migration finishes before Q-Day")
    _img_full(s, charts["migration"])
    _notes(s, "Standard path: ~7 years for broad adoption (BIP co-author estimate).\n"
              "Urgent path requires pre-work NOW: standards, wallet readiness, exchange coordination.\n"
              "Side-by-side with GRI risk signal to visualise the race.")

    # ── 20  Legacy coins ─────────────────────────────────────────────────────
    s = slide_base(prs, "Problem 2: legacy coins",
                   "Do nothing = stolen.  Restrict = censorship.  Both bad.")
    _two_col(s,
        "Options",
        ["Steal: let CRQC spend them, preserve neutrality",
         "Freeze: lock vulnerable UTXOs preemptively",
         "Hourglass: rate-limit movement of exposed coins"],
        "Consequences",
        ["Steal: price shock, but censorship resistance intact",
         "Freeze: market crash prevented, legitimacy damaged",
         "Hourglass: middle path, still contentious"])
    _callout(s, "Likely outcome: chain split risk. External forces will exploit the division.", color=RED)
    _notes(s, "This is governance territory, not just engineering.\n"
              "Hourglass BIP: rate-limit quantum-exposed coin movement.\n"
              "If freeze happens too early, legitimate users harmed.\n"
              "If too late, market crash from quantum theft.")

    # ── 21  What can you do ──────────────────────────────────────────────────
    s = slide_base(prs, "What can you do?")
    _two_col(s,
        "Today",
        ["Don't reuse addresses",
         "Avoid P2TR/P2PK for large holdings",
         "Don't share your xpub unnecessarily",
         "  (SMSF auditing, public Electrum nodes)",
         "Satoshi's shield: game-theory deterrent"],
        "When standards land",
        ["Migrate to P2MR when wallets support it",
         "Adopt tested PQC implementations",
         "Run a node to support the upgrade",
         "Fund dev grants and crypto research",
         "Participate in the consensus direction"])
    _notes(s, "Satoshi's shield: if Satoshi's coins are vulnerable, attackers face massive sell pressure.\n"
              "Game-theoretic deterrent against quantum theft.\n"
              "xpub exposure: SMSF auditing requirements and public Electrum nodes leak derivation paths.")

    # ── 22  Mythbusters ──────────────────────────────────────────────────────
    s = slide_base(prs, "Mythbusters")
    _two_col(s,
        "Myth",
        ["CRQC breaks all addresses instantly",
         "Grover kills SHA-256",
         "This is only a Bitcoin problem"],
        "Reality",
        ["Exposure type and time window matter",
         "Quadratic speedup shifts cost curves, not instant break",
         "ECDSA/RSA break is an internet-wide crisis"])
    _notes(s, "Calibration slide. Staged threats, not a single catastrophic event.\n"
              "Grover: 2^256 -> 2^128 search space. Still enormous.")

    # ── 23  Broader perspective ──────────────────────────────────────────────
    s = slide_base(prs, "CRQC most likely use cases",
                   "Bitcoin is one target among many")
    _bullets(s, [
        "SIGINT: crack the 'store now, decrypt later' vault",
        "Materials science and engineering simulation",
        "Pharmaceutical drug discovery",
        "Disease prevention and genomics research",
        "Physics: push the boundary of QFT",
    ], y=2.1, size=26)
    _notes(s, "Nation-state drivers keep quantum investment high regardless of Bitcoin.\n"
              "'Store now, decrypt later' is already happening — Bitcoin's pressure is real.")

    # ── 24  References ───────────────────────────────────────────────────────
    add_reference_slides(prs)

    # ── 25  Q&A ──────────────────────────────────────────────────────────────
    s = slide_base(prs, "Questions?",
                   "Where should Bitcoin move first?")
    _bullets(s, [
        "Priority 1:  stop exposing public keys",
        "Priority 2:  safe upgrade path (BIP-360)",
        "Priority 3:  social consensus on legacy coins",
    ], y=2.6, w=11.0, size=30)
    _notes(s, "Close with priorities. Invite involvement.\n"
              "Key question: what ships early without fragmenting consensus?")

    prs.save(output_pptx := OUT)
    print(f"Saved  {output_pptx}")
    return output_pptx


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("--no-odp", action="store_true", help="Skip ODP conversion")
    args = ap.parse_args()

    charts = generate_charts()
    pptx_path = build_deck(charts)

    if not args.no_odp:
        try:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "odp",
                 "--outdir", str(ROOT), str(pptx_path)],
                check=True, capture_output=True, timeout=60,
            )
            odp = pptx_path.with_suffix(".odp")
            print(f"Saved  {odp}")
        except (subprocess.CalledProcessError, FileNotFoundError) as e:
            print(f"ODP conversion skipped: {e}", file=sys.stderr)


if __name__ == "__main__":
    main()
